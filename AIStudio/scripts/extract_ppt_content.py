#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT内容提取脚本 - 用于公众号网文生成
提取PPT中的所有文本内容，包括标题、正文、表格等
"""

import os
import json
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_all_content_from_slide(slide):
    """
    从单个幻灯片中提取所有内容，包括文本、表格、图片描述等
    """
    slide_content = {
        "slide_number": slide.slide_id,
        "title": "",
        "text_content": [],
        "tables": [],
        "shapes": [],
        "notes": ""
    }
    
    # 提取幻灯片备注
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            slide_content["notes"] = notes_slide.notes_text_frame.text.strip()
    
    for shape in slide.shapes:
        shape_info = {
            "type": str(shape.shape_type),
            "text": "",
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
        }
        
        # 处理文本框
        if shape.has_text_frame:
            text_content = []
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ""
                for run in paragraph.runs:
                    paragraph_text += run.text
                if paragraph_text.strip():
                    text_content.append(paragraph_text.strip())
            
            shape_info["text"] = "\n".join(text_content)
            
            # 判断是否为标题（通常位置靠上且字体较大）
            if shape.top < 2000000:  # 位置靠上
                if not slide_content["title"]:
                    slide_content["title"] = shape_info["text"]
        
        # 处理表格
        elif shape.has_table:
            table_data = []
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell_idx, cell in enumerate(row.cells):
                    cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                    row_data.append(cell_text)
                table_data.append(row_data)
            
            shape_info["table_data"] = table_data
            slide_content["tables"].append({
                "data": table_data,
                "position": shape_info["position"]
            })
        
        # 处理图片（记录位置和大小信息）
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info["image_info"] = {
                "filename": getattr(shape.image, 'filename', 'unknown'),
                "size": f"{shape.width}x{shape.height}"
            }
        
        slide_content["shapes"].append(shape_info)
        
        # 收集所有文本内容
        if shape_info["text"]:
            slide_content["text_content"].append(shape_info["text"])
    
    return slide_content

def extract_ppt_structure(prs):
    """
    提取PPT的整体结构信息
    """
    structure = {
        "total_slides": len(prs.slides),
        "slide_layouts": [],
        "master_info": {
            "title": prs.core_properties.title or "未命名演示文稿",
            "author": prs.core_properties.author or "未知作者",
            "created": str(prs.core_properties.created) if prs.core_properties.created else "",
            "modified": str(prs.core_properties.modified) if prs.core_properties.modified else ""
        }
    }
    
    # 分析幻灯片布局
    for slide in prs.slides:
        layout_name = slide.slide_layout.name if slide.slide_layout else "未知布局"
        structure["slide_layouts"].append(layout_name)
    
    return structure

def main():
    parser = argparse.ArgumentParser(
        description="提取PPT内容用于AI生成公众号网文"
    )
    parser.add_argument("--input-ppt", required=True, help="输入的.pptx文件路径")
    parser.add_argument("--output-json", required=True, help="输出JSON文件路径")
    parser.add_argument("--extract-images", action="store_true", help="是否提取图片信息")
    
    args = parser.parse_args()
    
    print("--- Starting PPT Content Extraction for Article Generation ---")
    print(f"Input PPT: {args.input_ppt}")
    print(f"Output JSON: {args.output_json}")
    
    try:
        prs = Presentation(args.input_ppt)
        
        # 提取整体结构
        structure = extract_ppt_structure(prs)
        
        # 提取每页内容
        slides_content = []
        for i, slide in enumerate(prs.slides):
            print(f"  - Processing slide {i+1}/{len(prs.slides)}...")
            slide_content = extract_all_content_from_slide(slide)
            slide_content["slide_number"] = i + 1
            slides_content.append(slide_content)
        
        # 生成摘要信息
        summary = {
            "total_slides": len(slides_content),
            "slides_with_content": len([s for s in slides_content if s["text_content"]]),
            "total_text_blocks": sum(len(s["text_content"]) for s in slides_content),
            "total_tables": sum(len(s["tables"]) for s in slides_content),
            "estimated_word_count": sum(len(" ".join(s["text_content"])) for s in slides_content)
        }
        
        # 组合最终输出
        output_data = {
            "structure": structure,
            "summary": summary,
            "slides": slides_content,
            "extraction_info": {
                "extract_images": args.extract_images,
                "extraction_time": str(os.popen('date').read().strip())
            }
        }
        
        # 确保输出目录存在
        output_dir = os.path.dirname(args.output_json)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        
        # 保存结果
        with open(args.output_json, 'w', encoding='utf-8') as f:
            json.dump(output_data, f, ensure_ascii=False, indent=2)
        
        print(f"\nSuccessfully extracted PPT content:")
        print(f"  - Total slides: {summary['total_slides']}")
        print(f"  - Slides with content: {summary['slides_with_content']}")
        print(f"  - Total text blocks: {summary['total_text_blocks']}")
        print(f"  - Total tables: {summary['total_tables']}")
        print(f"  - Estimated word count: {summary['estimated_word_count']}")
        print(f"  - Saved to: {args.output_json}")
        
    except Exception as e:
        print(f"Error processing PPT: {e}")
        return 1
    
    return 0

if __name__ == "__main__":
    exit(main())
