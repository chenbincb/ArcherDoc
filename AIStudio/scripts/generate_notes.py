import os
import sys
import json
import argparse
import re
import base64
import requests
from openai import OpenAI
from pptx import Presentation

# --- 【核心升级】增强的文本提取函数 ---
def extract_text_from_slide(slide, image_path=None):
    """
    融合模式：同时提取原始文本和图片视觉描述。
    原始文本保证技术术语的精确性，图片描述提供逻辑布局信息。
    """
    raw_texts = []
    # 提取文字和表格
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        cell_text = cell.text_frame.text.strip()
                        if cell_text: raw_texts.append(cell_text)
        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run_text = run.text.strip()
                    if run_text: raw_texts.append(run_text)
    
    extracted_text = "\n".join(raw_texts)
    
    # 尝试使用多模态模型识别图片
    visual_description = ""
    if image_path:
        print(f"    尝试获取幻灯片视觉布局描述...")
        visual_description = recognize_image_with_multimodal(image_path) or ""
        if visual_description:
            print(f"    ✓ 获取视觉描述成功")
    
    # 将两者融合，互为补充
    combined_content = f"【原始文本内容】:\n{extracted_text}\n\n【视觉布局描述】:\n{visual_description}"
    return combined_content


# --- 【新增】从job目录获取图片路径函数 ---
def get_image_path_from_job(input_ppt_path, slide_index):
    """
    从输入PPT路径解析job目录，构建对应的幻灯片图片路径
    图片路径格式：{job_dir}/slides/slide_{slide_index}.png
    """
    try:
        # 从input_ppt_path解析job目录
        # input_ppt_path 格式通常为：/home/n8n/AIStudio/jobs/{job_id}/input.pptx
        input_dir = os.path.dirname(input_ppt_path)
        image_dir = os.path.join(input_dir, "slides")
        image_path = os.path.join(image_dir, f"slide_{slide_index}.png")
        
        # 检查图片文件是否存在
        if os.path.exists(image_path):
            return image_path
        return None
    except Exception as e:
        print(f"    Warning: 无法获取图片路径 - {e}")
        return None

# --- 【新增】使用多模态模型识别图片内容函数 ---
def recognize_image_with_multimodal(image_path):
    """
    使用本地vllm部署的Qwen3-VL-4B-Instruct模型识别图片内容
    返回图片的文本描述
    """
    try:
        # 硬编码模型参数
        model_config = {
            "baseUrl": "http://178.109.129.11:8008/v1",
            "model": "/home/n8n/Qwen3-VL/Qwen3-VL-4B-Instruct",
            "apiKey": "EMPTY"
        }
        
        # 读取图片二进制数据
        with open(image_path, "rb") as f:
            image_data = f.read()
        
        # 将图片转换为base64编码
        image_base64 = base64.b64encode(image_data).decode("utf-8")
        
        # 构建请求体
        payload = {
            "model": model_config["model"],
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "text", "text": "请详细描述这张幻灯片的内容，包括标题、正文、图表、图片等所有元素。输出简洁明了，直接给出描述结果。"},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_base64}"}}
                ]
            }],
            "temperature": 0.1
        }
        
        # 调用vllm API
        headers = {
            "Authorization": f"Bearer {model_config['apiKey']}",
            "Content-Type": "application/json"
        }
        
        response = requests.post(
            f"{model_config['baseUrl']}/chat/completions",
            headers=headers,
            json=payload,
            timeout=30  # 设置30秒超时
        )
        
        # 检查响应状态
        response.raise_for_status()
        
        # 解析响应结果
        result = response.json()
        if result and "choices" in result and result["choices"]:
            return result["choices"][0]["message"]["content"].strip()
        return None
    except Exception as e:
        print(f"    Warning: 多模态模型识别失败 - {e}")
        return None

# --- 【新增】提取PPT备注函数 ---
def extract_notes_from_slide(slide):
    """
    从PPT幻灯片中提取备注文字
    备注是演讲者在PPT中添加的演讲提示，通常在"备注"视图中编写
    """
    try:
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            return notes_text
        return ""
    except Exception as e:
        # 如果提取备注出错，返回空字符串，不影响主流程
        print(f"    Warning: 无法提取备注 - {e}")
        return ""


def main():
    """主执行函数，负责解析参数、调用AI、输出结果。"""

    parser = argparse.ArgumentParser(
        description="为PPT演示文稿生成讲稿。兼容OpenAI, OpenRouter及本地Ollama等服务。"
    )
    parser.add_argument("--input-ppt", required=True, help="输入的.pptx文件的相对路径。")
    parser.add_argument("--output-json", required=True, help="用于保存输出.json文件的相对路径。")
    parser.add_argument("--type", choices=["video", "image"], default="video", help="生成内容的类型：video=生成演讲稿，image=生成图片描述和提示词")

    # 从命令行参数获取AI配置，设置默认值
    parser.add_argument("--api-base-url", default="https://openrouter.ai/api/v1", help="大语言模型API的服务器地址。")
    parser.add_argument("--model-name", default="qwen/qwen3-235b-a22b:free", help="要使用的大语言模型的名称。")
    parser.add_argument("--api-key", default="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", help="用于访问API服务的密钥。")
    
    args = parser.parse_args()

    print("--- Starting Script: 1_generate_notes.py (Enhanced Content Generation) ---")
    print(f"Input PPT: {args.input_ppt}")
    print(f"Generation Type: {args.type}")
    print(f"Using Model: {args.model_name}")
    print(f"API Base URL: {args.api_base_url}")

    if not args.api_key or "xxxxxxxx" in args.api_key:
        print("Error: API key is not configured. Please provide a valid API key using --api-key parameter.")
        sys.exit(1)

    try:
        client = OpenAI(
            base_url=args.api_base_url,
            api_key=args.api_key
        )
    except Exception as e:
        print(f"Error: Failed to initialize OpenAI client. Details: {e}")
        return

    try:
        prs = Presentation(args.input_ppt)
        all_content = []

        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            print(f"  - Processing slide {slide_number}/{len(prs.slides)}...")

            # 获取当前幻灯片对应的图片路径
            image_path = get_image_path_from_job(args.input_ppt, i)

            # 使用增强的文本提取函数，优先使用图片识别
            slide_content = extract_text_from_slide(slide, image_path)
            notes_text = extract_notes_from_slide(slide)

            # 提取幻灯片标题（优先从PPT内容中获取第一行作为标题，如果没有内容则使用备注或默认标题）
            slide_lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
            if slide_lines:
                title = slide_lines[0]  # 使用PPT内容的第一行作为标题
            elif notes_text:
                title = notes_text.split('\n')[0]  # 使用备注的第一行作为标题
            else:
                title = f"幻灯片 {slide_number}"  # 默认标题

            if args.type == "video":
                # 视频模式：生成演讲稿
                if notes_text:
                    generated_note = notes_text
                    print(f"    ✓ 使用PPT备注 ({len(notes_text)}字符)")
                else:
                    if not slide_content.strip():
                        print(f"    ! 无文本内容且无备注")
                        generated_note = "（此页未识别到文本内容，请根据图片或图表自行描述。）"
                    else:
                        print(f"    AI生成讲稿中...")
                        prompt = f"""
                        你是一位顶级的演讲文稿撰写专家。
                        这是PPT第 {slide_number} 页上的所有文字内容，由不同文本框和表格单元格拼接而成：
                        ---
                        {slide_content}
                        ---
                        请根据这些内容，为这一页生成一段大约100-150字的、自然流畅、专业且引人入胜的演讲稿（内容较多的字数最多可以到300字）。
                        请直接输出演讲稿文本，不要包含"好的，这是您的演讲稿："等多余的前言或结语。
                        """
                        try:
                            response = client.chat.completions.create(
                                model=args.model_name,
                                messages=[{"role": "user", "content": prompt}]
                            )
                            generated_note = response.choices[0].message.content.strip()
                            print(f"    ✓ AI生成完成")
                        except Exception as e:
                            print(f"    ✗ AI生成失败: {e}")
                            generated_note = f"// AI生成失败: {e} //"

                all_content.append({"id": i, "note": generated_note})

            elif args.type == "image":
                # 图片模式：生成图片描述和提示词
                if not slide_content.strip():
                    print(f"    ! 无文本内容，使用默认描述")
                    generated_description = f"第{slide_number}页PPT内容：{title}"
                    suggested_prompt = f"关于 {title} 的逻辑图表, 科技风格, 结构化信息图表, 专业产品文档插图, 扁平化设计, 几何构图, 清晰的逻辑线条, 商务色调, 适合PPT展示, 无文字标签"
                    print(f"    ✓ 使用默认描述和提示词")
                else:
                    print(f"    AI生成图片描述和提示词中...")
                    
                    # 定义不同类型的专属强指令 (与前端 ImageReviewPage.tsx 保持完全一致)
                    type_instructions = """
1. **逻辑架构图 (Logical Architecture)**
   - 构图：2.5D等轴测 (Isometric View)，模块化堆叠。
   - 适用：系统分层、模块组成、架构设计。
   - 视觉：底部基础设施，中间平台服务，顶部应用场景。

2. **业务流程图 (Business Process)**
   - 构图：2D扁平化，严格从左到右 (Left-to-Right)。
   - 适用：操作步骤、数据流转、业务闭环。
   - 视觉：输入源 -> 处理引擎 -> 输出结果，用箭头连接。

3. **网络拓扑图 (Network Topology)**
   - 构图：广角俯视 (Top-down)，星系分布或网状。
   - 适用：节点部署、多地多中心、互联互通。
   - 视觉：核心节点向周边辐射，强调连接线。

4. **数据可视化 (Data Visualization)**
   - 构图：正视UI界面 (Screen Mockup)，仪表盘布局。
   - 适用：数据统计、趋势分析、监控大屏。
   - 视觉：柱状图、折线图、KPI卡片，不要画实物。

5. **产品路线图 (Roadmap)**
   - 构图：2D水平时间轴 (Timeline)。
   - 适用：版本规划、发展历程、演进阶段。
   - 视觉：主轴线上分布里程碑节点。

6. **封面/通用页 (Cover/General)**
   - 构图：极简留白，抽象几何背景，中心化排版。
   - 适用：封面、目录、过渡页、纯文字总结。
   - 视觉：大气的主题背景 (Key Visual)，品牌色光影，无具体技术细节。
"""

                    prompt = f"""
你是一位专注【私有云/B端软件产品】的资深技术分析师和视觉设计师。
你的任务是将PPT内容转化为**深度理解后的技术描述**和**结构化图解提示词**。

<slide_content>
<title>{title}</title>
<content>{slide_content}</content>
</slide_content>

<task>
【步骤 1：判断页面性质与内容理解】
请先判断这张PPT的性质（是封面？目录？还是正文？）。
- **如果是封面/目录/过渡页**：请侧重描述**视觉氛围**和**品牌调性**。严禁脑补具体的技术架构细节！不要因为标题里有"存算分离"就去画存储架构图，这只是一张封面。
- **如果是正文内容页**：请像分析师一样拆解逻辑，识别技术实体（组件）、逻辑行为（关系）和核心诉求（价值）。

【步骤 2：智能分类】
根据页面性质，从以下6种类型中选择最匹配的一种：
{type_instructions}

【步骤 3：生成结构化提示词】
基于你的深度理解，进行视觉建模，严格执行以下要求。
</task>

<design_guidelines>
<composition_principles>
- 根据内容自动设计最完美的构图
- 重点突出核心概念，避免信息过载
- 使用装饰性元素填补空白，保持画面平衡
- 避免过度拥挤或过度留白
</composition_principles>

<text_rendering_rules>
【核心原则】
- 如需渲染文字，不重不漏地包含所有关键信息
- 保持原文的逻辑层次和重点强调

【格式规范】
- 禁止使用markdown格式符号（如 # * - 等）
- 标题使用字号和粗细区分，不添加符号
- 列表项使用缩进组织，不添加项目符号

【内容限制】
- 保留技术缩写的英文形式（API、HTTP、JSON、Cloud、DB、SaaS、PaaS、IaaS等）
- 其他标签和说明文字使用中文
- 如果无法保证汉字清晰，生成空白文本框，不要生成乱码英文
</text_rendering_rules>

<quality_standards>
- 视觉重心突出，主体明确
- 元素分布均衡，有呼吸感
- 引导线清晰，逻辑流畅
- 符合阅读习惯（从左到右，从上到下）
- 专业商务PPT风格，简洁现代
</quality_standards>
</design_guidelines>

<output_format>
请返回一个标准的JSON对象：
{{
  "description": "深度技术转译结果。如果是封面，请描述'大气、专业的开场视觉'；如果是正文，请详细描述业务逻辑",
  "identifiedType": "从上述6种类型中选择其一（如：封面/通用页）",
  "suggestedPrompt": "结构化提示词，包含：1.场景构图 2.核心组件 3.逻辑交互 4.文本标签 5.视觉风格"
}}
"""

                    try:
                        response = client.chat.completions.create(
                            model=args.model_name,
                            messages=[{"role": "user", "content": prompt}]
                        )
                        ai_result = response.choices[0].message.content.strip()

                        # 尝试解析JSON
                        try:
                            # 查找JSON内容
                            json_match = re.search(r'\{.*\}', ai_result, re.DOTALL)
                            if json_match:
                                result = json.loads(json_match.group())
                                generated_description = result.get("description", f"第{slide_number}页PPT内容：{title}。{slide_content[:100]}...")
                                suggested_prompt = result.get("suggestedPrompt", f"关于 {title} 的逻辑图表, 科技风格, 结构化信息图表")
                                identified_type = result.get("identifiedType", "自动识别")
                                print(f"    ✓ AI生成完成，识别类型：{identified_type}")
                            else:
                                # JSON解析失败，使用默认值
                                generated_description = f"第{slide_number}页PPT内容：{title}。{slide_content[:100]}..."
                                suggested_prompt = f"关于 {title} 的逻辑图表, 科技风格, 结构化信息图表, 专业产品文档插图, 扁平化设计, 几何构图, 清晰的逻辑线条, 商务色调, 适合PPT展示, 无文字标签"
                                print(f"    ⚠️ JSON解析失败，使用默认提示词")
                        except json.JSONDecodeError as je:
                            # JSON解析失败，使用默认值
                            generated_description = f"第{slide_number}页PPT内容：{title}。{slide_content[:100]}..."
                            suggested_prompt = f"关于 {title} 的逻辑图表, 科技风格, 结构化信息图表, 专业产品文档插图, 扁平化设计, 几何构图, 清晰的逻辑线条, 商务色调, 适合PPT展示, 无文字标签"
                            print(f"    ⚠️ JSON解析失败: {je}，使用默认提示词")

                    except Exception as e:
                        print(f"    ✗ AI生成失败: {e}")
                        generated_description = f"第{slide_number}页PPT内容：{title}。{slide_content[:100]}..."
                        suggested_prompt = f"关于 {title} 的逻辑图表, 科技风格, 结构化信息图表, 专业产品文档插图, 扁平化设计, 几何构图, 清晰的逻辑线条, 商务色调, 适合PPT展示, 无文字标签"

                all_content.append({
                    "id": i,
                    "title": title,
                    "content": slide_content,
                    "description": generated_description,
                    "suggestedPrompt": suggested_prompt
                })

        output_dir = os.path.dirname(args.output_json)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
            
        with open(args.output_json, 'w', encoding='utf-8') as f:
            json.dump(all_content, f, ensure_ascii=False, indent=4)

        print(f"\nSuccessfully generated notes and saved to: {args.output_json}")

    except Exception as e:
        print(f"An error occurred during presentation processing: {e}")

if __name__ == "__main__":
    main()